#!/usr/bin/env python3
# Build a native, editable PPTX of the 「舊地圖散策」project deck.
# Mirrors slides.html (vintage paper palette). Run: python3 build_deck.py
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── palette (from slides.html) ───────────────────────────
PAPER   = RGBColor(0xF1, 0xEF, 0xE9)
PAPER2  = RGBColor(0xEA, 0xE8, 0xE2)
INK     = RGBColor(0x1F, 0x1D, 0x19)
SUMI    = RGBColor(0x14, 0x13, 0x0F)
GRAPH   = RGBColor(0x4A, 0x46, 0x3E)
ASH     = RGBColor(0x6E, 0x6A, 0x60)
SILVER  = RGBColor(0x9C, 0x98, 0x8C)
MIST    = RGBColor(0xD8, 0xD4, 0xCA)
STONE   = RGBColor(0xC9, 0xC4, 0xB8)
INDIGO  = RGBColor(0x2E, 0x43, 0x74)
VERM    = RGBColor(0xD6, 0x5B, 0x3A)

ZH_SERIF = "Noto Serif TC"      # falls back gracefully if absent
ZH_SANS  = "Noto Sans TC"
EN_SERIF = "EB Garamond"
MONO     = "JetBrains Mono"

EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)   # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def inch(v): return Emu(int(v * EMU))

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    sp = bg._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _set_font(run, font, size, color, bold=False, italic=False, spacing=None):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    # set east-asian font too
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)
    if spacing is not None:
        rPr.set('spc', str(int(spacing * 100)))  # pt → 1/100 pt units used by spc

def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p

def run(p, text, font, size, color, bold=False, italic=False, spacing=None):
    r = p.add_run(); r.text = text
    _set_font(r, font, size, color, bold, italic, spacing)
    return r

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def accent_bar(s, x, y, w=0.85, color=VERM):
    rect(s, x, y, w, 0.03, fill=color)

def kicker(s, text, x=0.85, y=0.7):
    accent_bar(s, x, y + 0.07, 0.55, VERM)
    tf = textbox(s, x + 0.7, y, 8, 0.3)
    run(para(tf, True), text, MONO, 11, VERM, spacing=2.4)

def pageno(s, n):
    tf = textbox(s, 11.7, 6.9, 1.4, 0.3)
    p = para(tf, True); p.alignment = PP_ALIGN.RIGHT
    run(p, f"{n:02d} / 10", MONO, 9, SILVER, spacing=2.0)

def brandfoot(s, text="Project Deck"):
    tf = textbox(s, 0.85, 6.9, 4, 0.3)
    run(para(tf, True), text.upper(), MONO, 8, SILVER, spacing=2.0)

def title_block(s, zh, en, x=0.85, y=1.05):
    tf = textbox(s, x, y, 11.5, 1.5)
    p = para(tf, True)
    run(p, zh, ZH_SERIF, 30, SUMI, bold=True, spacing=1.0)
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    run(p2, en, EN_SERIF, 14, INDIGO, italic=True)

def card(s, x, y, w, h, title_zh, title_en, desc, icon="", accent=INDIGO):
    box = rect(s, x, y, w, h, fill=PAPER2, line=STONE, line_w=1.0)
    # left accent stripe
    rect(s, x, y, 0.045, h, fill=accent)
    tf = textbox(s, x + 0.28, y + 0.2, w - 0.5, h - 0.36, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, True)
    if icon: run(p, icon + "  ", ZH_SANS, 15, INK)
    run(p, title_zh, ZH_SERIF, 15, INK, bold=True, spacing=0.6)
    if title_en:
        run(p, "  " + title_en, EN_SERIF, 10.5, INDIGO, italic=True)
    if desc:
        pd = tf.add_paragraph(); pd.space_before = Pt(5)
        run(pd, desc, ZH_SANS, 11, ASH)
    return box

# ════════════════════════════════════════════════════════
# 01 — Title
# ════════════════════════════════════════════════════════
s = slide()
# stamp
stamp = rect(s, 0.85, 1.5, 3.2, 0.42, fill=None, line=VERM, line_w=1.25)
tf = textbox(s, 0.85, 1.5, 3.2, 0.42, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
run(p, "TAIWAN · OLD-MAP STROLL", MONO, 10, VERM, spacing=2.6)
# big title
tf = textbox(s, 0.8, 2.25, 11.5, 1.7)
run(para(tf, True), "舊地圖散策", ZH_SERIF, 60, SUMI, bold=True, spacing=3.0)
tf = textbox(s, 0.85, 3.95, 11.5, 0.8)
run(para(tf, True), "Walk the island through a century of maps",
    EN_SERIF, 22, INDIGO, italic=True)
# year line
tf = textbox(s, 0.85, 4.75, 11.5, 0.4)
p = para(tf, True)
run(p, "1897    ——    1944    ——    2003    ——    今 NOW", MONO, 12, GRAPH, spacing=1.5)
# meta
tf = textbox(s, 0.85, 5.35, 11.5, 1.3)
for i, ln in enumerate([
    "把百年來的臺灣歷史地圖，疊回你腳下的街道",
    "A walking-first PPGIS · PWA · iOS / Android · 中 / EN / 日",
    "吳昀慶、張芸翠、朱穎芃",
]):
    p = para(tf, i == 0); p.space_after = Pt(4)
    run(p, ln, MONO if i else ZH_SANS, 11 if i else 12, ASH if i else GRAPH, spacing=1.2)
brandfoot(s, "Project Deck"); pageno(s, 1)

# ════════════════════════════════════════════════════════
# 02 — The idea
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "01 · THE IDEA")
title_block(s, "你腳下的馬路，百年前是什麼？", "A time machine you operate with your feet")
# lead + ticks (left)
tf = textbox(s, 0.85, 2.7, 7.4, 1.3)
run(para(tf, True),
    "舊地圖散策把 1897 年以來的臺灣歷史地圖，疊在今日的街道上。"
    "你現在走的這條路，可能曾是水圳、糖鐵，或一座消失的聚落。",
    ZH_SERIF, 14.5, GRAPH, spacing=0.5)
ticks = [
    ("不只是看圖", "帶著手機走出門，用 GPS 讓歷史在你腳下「對位」"),
    ("不只是懷舊", "把分散在各機構的圖資，整理成人人可玩的介面"),
    ("邊走邊收集", "走到歷史現場，老明信片就會「沖洗」顯影"),
]
tf = textbox(s, 0.85, 4.15, 7.4, 2.4)
for i, (b, d) in enumerate(ticks):
    p = para(tf, i == 0); p.space_after = Pt(9)
    run(p, "—  ", MONO, 13, VERM)
    run(p, b + "　", ZH_SERIF, 13.5, INK, bold=True)
    run(p, d, ZH_SANS, 12.5, GRAPH)
# right note card
card(s, 8.7, 2.8, 3.75, 3.4,
     "核心命題", "", "", icon="🗺️", accent=VERM)
tf = textbox(s, 9.0, 3.55, 3.2, 2.5, anchor=MSO_ANCHOR.TOP)
for i, ln in enumerate(["歷史不在課本裡，", "就在你走的這條路上。", "",
                        "讓「地點」成為時間的入口，", "讓散步成為一場考古。"]):
    p = para(tf, i == 0); p.space_after = Pt(3)
    run(p, ln, ZH_SERIF, 13, GRAPH, spacing=0.4)
pageno(s, 2)

# ════════════════════════════════════════════════════════
# 03 — Content & data
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "02 · CONTENT & DATA")
title_block(s, "內容與資料", "23 historical layers · 22 cities · 200+ postcards")
stats = [("23", "歷史圖層 LAYERS"), ("22", "縣市 CITIES"),
         ("200+", "老明信片 POSTCARDS"), ("3", "語言 中 / EN / 日")]
x = 0.85
for n, l in stats:
    tf = textbox(s, x, 2.65, 2.9, 1.4)
    run(para(tf, True), n, EN_SERIF, 50, VERM)
    p = tf.add_paragraph(); p.space_before = Pt(2)
    run(p, l, MONO, 8.5, ASH, spacing=1.6)
    x += 3.0
# timeline (4 eras)
eras = [
    ("1897 – 1944", "日治時期", "臺灣堡圖、蕃地地形圖、二萬五千分一地形圖…（10 圖層）"),
    ("1944 – 1945", "美軍時期", "U.S. Army 地形圖與城市地圖（4 圖層）"),
    ("1950 – 2003", "戰後經建版", "臺灣地形圖、經建版一～四版（7 圖層）"),
    ("1966 / 1969", "CORONA 衛星", "冷戰時期解密衛星影像（2 圖層）"),
]
rect(s, 0.85, 4.55, 11.6, 0.012, fill=STONE)
x = 0.85; w = 2.9
for i, (yr, nm, sub) in enumerate(eras):
    tf = textbox(s, x, 4.75, w - 0.15, 1.5)
    run(para(tf, True), yr, MONO, 10, VERM, spacing=1.0)
    pn = tf.add_paragraph(); pn.space_before = Pt(3)
    run(pn, nm, ZH_SERIF, 12.5, INK, bold=True)
    ps = tf.add_paragraph(); ps.space_before = Pt(3)
    run(ps, sub, ZH_SANS, 9.5, ASH)
    x += w
# sources
rect(s, 0.85, 6.45, 11.6, 0.012, fill=STONE)
tf = textbox(s, 0.85, 6.55, 11.6, 0.7)
p = para(tf, True)
run(p, "資料來源  ", MONO, 9, GRAPH, bold=True, spacing=1.2)
run(p, "歷史圖資 · 中央研究院 GIS   |   現代底圖 · CARTO + OpenStreetMap   |   "
       "正射影像 · 國土測繪中心 NLSC   |   日本圖資 · 農研機構 / 國土地理院",
    MONO, 9, ASH, spacing=0.6)
pageno(s, 3)

# ════════════════════════════════════════════════════════
# 04 — Core features (2×2)
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "03 · CORE FEATURES")
title_block(s, "讀圖的四種方式", "Four ways to read time on the map")
feats = [
    ("圖層疊加", "Overlay & opacity", "任選一張古地圖疊上現代街道，用透明度滑桿細調，讓今昔同時可見。", "🎚️", INDIGO),
    ("對照拉條", "Swipe compare", "左右滑動分界線，一邊現代、一邊古圖；也能古圖對古圖，比較年代。", "↔️", INDIGO),
    ("四代對照", "Four eras", "四個年代並排同步縮放平移，一眼看完城市百年的長成過程。", "⊞", INDIGO),
    ("時光快轉", "Time-lapse", "自選年代依序淡入淡出，從 1897 一路播到今天，像會動的地圖。", "▶", VERM),
]
gx, gy, gw, gh, gap = 0.85, 2.75, 5.7, 1.75, 0.25
for i, (zh, en, d, ic, ac) in enumerate(feats):
    cx = gx + (i % 2) * (gw + gap)
    cy = gy + (i // 2) * (gh + gap)
    card(s, cx, cy, gw, gh, zh, en, d, icon=ic, accent=ac)
pageno(s, 4)

# ════════════════════════════════════════════════════════
# 05 — Explore & collect (2×3)
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "04 · EXPLORE & COLLECT")
title_block(s, "把散步變成一場考古", "Gamified field exploration")
items = [
    ("明信片顯影", "走到歷史現場 50 公尺內，日治老明信片自動「沖洗」收進圖鑑。", "📮", VERM),
    ("每日抽卡", "每 4 小時、每日上限 3 張，沒走到也能先收藏回味。", "🎴", INDIGO),
    ("附近的時光", "一開 App 就推鄰近現場故事卡，連偏鄉也有探索方向。", "🧭", INDIGO),
    ("散策記錄", "GPS 記錄路線、距離與時間，沿途拍照釘在地圖上。", "👣", INDIGO),
    ("縣市印章", "走訪 22 縣市收集印章，集成一本旅行護照。", "🪧", INDIGO),
    ("散策月曆", "把每天的足跡標進月曆，回顧自己的時間旅行。", "📅", INDIGO),
]
gx, gy, gw, gh, gap = 0.85, 2.75, 3.73, 1.75, 0.2
for i, (zh, d, ic, ac) in enumerate(items):
    cx = gx + (i % 3) * (gw + gap)
    cy = gy + (i // 3) * (gh + gap)
    card(s, cx, cy, gw, gh, zh, "", d, icon=ic, accent=ac)
pageno(s, 5)

# ════════════════════════════════════════════════════════
# 06 — Architecture (3 rows)
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "05 · ARCHITECTURE")
title_block(s, "基本架構", "A pure front-end PWA — no server required")
rows = [
    ("SHELL", [
        ("單檔 PWA", "HTML + CSS + Vanilla JS", INDIGO),
        ("Leaflet.js", "地圖 / 圖磚 · 淡入淡出", None),
        ("Web Audio API", "即時合成音效，無音檔", None),
    ]),
    ("OFFLINE", [
        ("Service Worker", "App Shell 預快取 · 離線可用", VERM),
        ("圖磚 LRU 快取", "~1200 張上限 · 自動驅逐", None),
        ("localStorage", "收藏 / 路線 / 設定存本機", None),
    ]),
    ("DATA / EDGE", [
        ("Cloudflare Workers", "中研院圖磚代理 + CORS", INDIGO),
        ("中研院 / NLSC / 日本", "歷史圖磚與正射影像", None),
        ("Capacitor 7", "同碼上架 iOS / Android", None),
    ]),
]
ry = 2.75
for tag, boxes in rows:
    tf = textbox(s, 0.85, ry + 0.1, 1.7, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, True); p.alignment = PP_ALIGN.RIGHT
    run(p, tag, MONO, 9, ASH, spacing=1.6)
    bx = 2.75; bw = 3.05; bgap = 0.22
    for b, sub, ac in boxes:
        box = rect(s, bx, ry, bw, 0.95, fill=PAPER2, line=(ac or STONE), line_w=1.25)
        if ac: rect(s, bx, ry, 0.045, 0.95, fill=ac)
        tf = textbox(s, bx + 0.25, ry, bw - 0.4, 0.95, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True), b, ZH_SERIF, 12.5, INK, bold=True, spacing=0.4)
        ps = tf.add_paragraph(); ps.space_before = Pt(2)
        run(ps, sub, MONO, 8.5, ASH, spacing=0.4)
        bx += bw + bgap
    if tag != "DATA / EDGE":
        tf = textbox(s, 6.0, ry + 0.97, 1, 0.3)
        p = para(tf, True); p.alignment = PP_ALIGN.CENTER
        run(p, "▼", ZH_SANS, 12, SILVER)
    ry += 1.42
pageno(s, 6)

# ════════════════════════════════════════════════════════
# 07 — Offline & performance
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "06 · OFFLINE & PERFORMANCE")
title_block(s, "離線優先，走到山裡也能用", "Built to work without a signal")
ticks = [
    ("分桶快取策略", "App Shell、執行時資源、圖磚、離線包各自獨立，升級不互相清空"),
    ("圖磚 LRU 驅逐", "常用圖磚留下、舊的自動淘汰，控制儲存在 100–150 MB 內"),
    ("離線包下載", "把當前範圍整包釘住，出發前載好，山區無訊號照樣疊圖"),
    ("網路優先 / 快取優先", "HTML 走網路優先即時更新；圖磚走快取優先省流量"),
]
tf = textbox(s, 0.85, 2.8, 7.2, 3.3)
for i, (b, d) in enumerate(ticks):
    p = para(tf, i == 0); p.space_after = Pt(13)
    run(p, "—  ", MONO, 13, VERM)
    run(p, b, ZH_SERIF, 13.5, INK, bold=True)
    pd = tf.add_paragraph(); pd.space_after = Pt(13)
    run(pd, "     " + d, ZH_SANS, 12, GRAPH)
# SW buckets card
card(s, 8.5, 2.85, 3.95, 3.1, "SW 快取分桶", "", "", icon="⚙️", accent=VERM)
tf = textbox(s, 8.85, 3.65, 3.3, 2.2)
for i, ln in enumerate([
    "shell-cache　·　App 外殼",
    "runtime-cache　·　CDN / 字型",
    "tiles-cache　·　圖磚 (LRU 1200)",
    "pinned-cache　·　離線包 (永不驅逐)",
]):
    p = para(tf, i == 0); p.space_after = Pt(9)
    run(p, ln, MONO, 10.5, GRAPH, spacing=0.4)
pageno(s, 7)

# ════════════════════════════════════════════════════════
# 08 — Creative craft (2×3)
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "07 · CREATIVE CRAFT")
title_block(s, "讓人想一直走下去的小巧思", "The details that make it feel alive")
items = [
    ("走到才「顯影」", "未抵達的明信片以模糊泛黃預覽呈現，親自走到才全彩沖洗 — 保留抵達的獎勵感。", "🎞️", VERM),
    ("純合成音效", "翻頁、顯影、抽到稀有卡的聲音全由 Web Audio 即時合成，零音檔。", "🔊", INDIGO),
    ("老紙質美學", "紙張底色、噪點顆粒、襯線字與朱紅印記，介面像一份會動的舊地圖。", "📜", INDIGO),
    ("拖曳式抽屜", "底部面板 1:1 跟手拖曳、依速度甩動吸附；並修掉 Android 合成殘影。", "↕️", INDIGO),
    ("淡入淡出疊圖", "切換年代時圖磚交叉淡化，像底片疊印，而非生硬的瞬間跳換。", "✨", INDIGO),
    ("漸進式引導", "首訪導覽、今日明信片頁籤、可收合面板，讓複雜功能慢慢被發現。", "🧑‍🏫", INDIGO),
]
gx, gy, gw, gh, gap = 0.85, 2.75, 3.73, 1.78, 0.2
for i, (zh, d, ic, ac) in enumerate(items):
    cx = gx + (i % 3) * (gw + gap)
    cy = gy + (i // 3) * (gh + gap)
    card(s, cx, cy, gw, gh, zh, "", d, icon=ic, accent=ac)
pageno(s, 8)

# ════════════════════════════════════════════════════════
# 09 — Research framing
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "08 · RESEARCH FRAMING")
title_block(s, "一個「步行核心」的公眾參與式 GIS", "Walking-first PPGIS — and where it's heading")
tf = textbox(s, 0.85, 2.75, 7.2, 1.5)
p = para(tf, True)
run(p, "多數 PPGIS 讓人坐在電腦前點圖填報；本專案把參與的觸發點移到", ZH_SANS, 12.5, GRAPH)
run(p, "「身體實際走到的位置」", ZH_SANS, 12.5, INK, bold=True)
run(p, "，以接近觸發的收集與在地敘事，重新定義公眾如何與歷史地理互動。",
    ZH_SANS, 12.5, GRAPH)
tf = textbox(s, 0.85, 4.4, 7.2, 2.0)
for i, (b, d) in enumerate([
    ("現況", "非同步勘誤閉環：使用者回報釘位 / 史料 → 人工審核 → 釋出"),
    ("下一步", "即時 PPGIS：多人同步的知識共創"),
]):
    p = para(tf, i == 0); p.space_after = Pt(12)
    run(p, "—  ", MONO, 13, VERM)
    run(p, b + "　", ZH_SERIF, 13.5, INK, bold=True)
    run(p, d, ZH_SANS, 12, GRAPH)
# roadmap card
card(s, 8.5, 2.85, 3.95, 3.3, "未來架構 Roadmap", "", "", icon="🚀", accent=INDIGO)
tf = textbox(s, 8.85, 3.65, 3.35, 2.4)
p = para(tf, True)
run(p, "Supabase", ZH_SERIF, 12.5, INK, bold=True)
run(p, " (PostgreSQL + PostGIS + Realtime)", MONO, 9, ASH)
for ln in ["· 使用者勘誤上傳", "· 多源交叉驗證 / 研究者審核", "· 通過後即時推播至所有用戶端"]:
    pp = tf.add_paragraph(); pp.space_before = Pt(7)
    run(pp, ln, ZH_SANS, 11.5, GRAPH)
pageno(s, 9)

# ════════════════════════════════════════════════════════
# 10 — Closing
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "10 · CLOSING")
tf = textbox(s, 0.85, 2.4, 11.5, 2.0)
p = para(tf, True)
run(p, "歷史不在課本裡，", ZH_SERIF, 38, SUMI, bold=True, spacing=1.5)
p2 = tf.add_paragraph(); p2.space_before = Pt(8)
run(p2, "就在你", ZH_SERIF, 38, SUMI, bold=True, spacing=1.5)
run(p2, "走的這條路", ZH_SERIF, 38, VERM, bold=True, spacing=1.5)
run(p2, "上。", ZH_SERIF, 38, SUMI, bold=True, spacing=1.5)
tf = textbox(s, 0.85, 5.05, 11.5, 1.6)
for i, ln in enumerate([
    "舊地圖散策 · Taiwan Old-Map Stroll",
    "PWA · iOS / Android · 三語 · 離線可用 · 開源",
    "圖資 © 中央研究院 GIS · NLSC · CARTO / OpenStreetMap",
    "吳昀慶、張芸翠、朱穎芃",
]):
    p = para(tf, i == 0); p.space_after = Pt(4)
    run(p, ln, MONO, 10.5, ASH, spacing=1.2)
brandfoot(s, "Thank you"); pageno(s, 10)

out = "舊地圖散策_專案簡報.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
